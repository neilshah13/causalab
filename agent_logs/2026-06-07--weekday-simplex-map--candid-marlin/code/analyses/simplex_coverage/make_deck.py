"""Build the final slide deck (result/slides/weekday_map_slides.pptx), 16:9.

Style follows the TARA example deck: white background, small red ALL-CAPS kicker, large bold
black title, figure-dominant body, bold one-line takeaway at the bottom. Assets from
make_slide_assets.py (+ G1 GIF, which animates in PowerPoint).

Run from session root: python3 code/analyses/simplex_coverage/make_deck.py
"""

import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image

S = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "..")
SL = os.path.join(S, "result", "slides")
RF = os.path.join(S, "result", "figures")
W = os.path.join(S, "artifacts", "weekday_simplex", "llama31_8b", "simplex_coverage")

RED = RGBColor(0xB0, 0x30, 0x30)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
PAGE_W, PAGE_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = PAGE_W
prs.slide_height = PAGE_H
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, size, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             font="Arial", italic=False, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font
        f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font; f.italic = italic
    return tb


def add_img_fit(slide, path, x, y, max_w, max_h, align="center"):
    iw, ih = Image.open(path).size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    ox = x + (max_w - w) // 2 if align == "center" else x
    oy = y + (max_h - h) // 2
    slide.shapes.add_picture(path, ox, oy, width=w, height=h)


def header(slide, kicker, title):
    add_text(slide, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.3),
             kicker, 13, bold=True, color=RED)
    add_text(slide, Inches(0.55), Inches(0.58), Inches(12.2), Inches(0.75),
             title, 30, bold=True)


def takeaway(slide, text, y=Inches(6.55)):
    add_text(slide, Inches(0.55), y, Inches(12.2), Inches(0.6), text, 14.5, bold=True,
             color=BLACK)


# ---------------------------------------------------------------- 0 title
s = prs.slides.add_slide(BLANK)
add_text(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.6),
         "The manifold isn't the map", 44, bold=True)
add_text(s, Inches(0.8), Inches(3.5), Inches(11.7), Inches(1.0),
         "Mapping the geometry between activation space and behaviour —\n"
         "extending Goodfire's manifold steering (days of the week, Llama-3.1-8B)",
         19, color=GRAY, line_spacing=1.15)
add_text(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
         "TARA — weekday-simplex-map session, June 2026", 12, color=GRAY)

# ---------------------------------------------------------------- 1 problem: centroids
s = prs.slides.add_slide(BLANK)
header(s, "THE PROBLEM", "The manifold approach is limited")
add_img_fit(s, os.path.join(SL, "S1a_manifold_is_centroids.png"),
            Inches(0.55), Inches(1.5), Inches(12.2), Inches(5.1))
takeaway(s, "The “manifold” is 7 averaged points plus a curve-fitting step — the space between "
            "the points is never measured, only assumed.")

# ---------------------------------------------------------------- 2 problem: vectors match
s = prs.slides.add_slide(BLANK)
header(s, "THE PROBLEM", "The curve isn't doing the work — its points are")
add_img_fit(s, os.path.join(SL, "S1b_vectors_match_manifold.png"),
            Inches(0.55), Inches(1.5), Inches(12.2), Inches(5.1))
takeaway(s, "Plain vectors from centroid to centroid steer exactly as well as the fitted manifold. "
            "The geometry is in how the days are ARRANGED — the fitted connections between them are not real.")

# ---------------------------------------------------------------- 3 method
s = prs.slides.add_slide(BLANK)
header(s, "OUR METHOD", "Map the whole region — don't fit a curve")
add_img_fit(s, os.path.join(SL, "S2_method_flow.png"),
            Inches(0.55), Inches(1.45), Inches(12.2), Inches(2.5))
add_img_fit(s, os.path.join(RF, "G1_discovering_the_map.gif"),
            Inches(0.55), Inches(3.95), Inches(12.2), Inches(2.55))
takeaway(s, "Every step is measured. The geometry of the map is discovered, not assumed.")

# ---------------------------------------------------------------- 4 result 1
s = prs.slides.add_slide(BLANK)
header(s, "RESULT 1", "Steering along the discovered map is more natural")
add_img_fit(s, os.path.join(SL, "S3_naturalness_bars.png"),
            Inches(0.4), Inches(1.5), Inches(6.0), Inches(5.0), align="left")
gen = json.load(open(os.path.join(W, "steer_gen_demo_L23.json")))["results"]
samples = [
    ("steered → Saturday", gen[1]["carrier"], gen[1]["Sat"]["text"], "Saturday"),
    ("steered → Sunday", gen[2]["carrier"], gen[2]["Sun"]["text"], "Sunday"),
    ("steered → Wednesday", gen[0]["carrier"], gen[0]["Wed"]["text"], "Wednesday"),
]
y = Inches(1.55)
add_text(s, Inches(6.7), y, Inches(6.1), Inches(0.35),
         "…and the model keeps talking like nothing happened (P(day) ≈ 0.99):", 13, bold=True)
y = Inches(2.0)
for label, carrier, text, day in samples:
    tb = s.shapes.add_textbox(Inches(6.7), y, Inches(6.1), Inches(1.45))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = label; r.font.size = Pt(11.5); r.font.bold = True
    r.font.color.rgb = RED; r.font.name = "Arial"
    p2 = tf.add_paragraph(); p2.line_spacing = 1.05
    full = (carrier + text).replace("<|end_of_text|>", "")
    pre, _, post = full.partition(day)
    for chunk, bold in ((pre, False), (day, True), (post[:120] + "…", False)):
        r = p2.add_run(); r.text = chunk
        r.font.size = Pt(11.5); r.font.bold = bold; r.font.italic = True
        r.font.color.rgb = BLACK if bold else GRAY; r.font.name = "Georgia"
    y = y + Inches(1.55)
takeaway(s, "Same start and end points, same prompts, the paper's own naturalness score: "
            "2–5× closer to natural behaviour — and the steered model writes coherent text.")

# ---------------------------------------------------------------- 5 result 2
s = prs.slides.add_slide(BLANK)
header(s, "RESULT 2", "Two days of prompts recover the whole week")
add_img_fit(s, os.path.join(SL, "S4_recover_all_days.png"),
            Inches(0.4), Inches(1.45), Inches(8.3), Inches(5.2), align="left")
add_text(s, Inches(9.0), Inches(1.9), Inches(3.8), Inches(4.4),
         "Build the map from Monday + Thursday anchors ONLY.\n\n"
         "Walking it reaches every other day (Tue, Wed, Fri, Sat, Sun) at P(day) up to 1.0, "
         "without a single prompt about them.\n\n"
         "Red circles mark where the deleted days' anchors really live — the walk finds "
         "exactly those regions.",
         14, color=GRAY, line_spacing=1.15)
takeaway(s, "Behaviour can be discovered without prompting for it — and if the map is measured well, "
            "we can steer strongly toward what we find.")

out = os.path.join(SL, "weekday_map_slides.pptx")
prs.save(out)
print("saved", out)
